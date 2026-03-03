import streamlit as st
from openai import OpenAI
from zhipuai import ZhipuAI
import docx
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import json
import io
import os
import requests
from duckduckgo_search import DDGS

# ==========================================
# 1. 页面与全局设置
# ==========================================
st.set_page_config(page_title="物理教研 Agent", page_icon="⚛️", layout="wide", initial_sidebar_state="expanded")

try:
    api_key_brain = st.secrets["DEEPSEEK_API_KEY"]
    api_key_paint = st.secrets["ZHIPU_API_KEY"]
except KeyError:
    st.error("⚠️ 未找到 API Key！请检查 .streamlit/secrets.toml 配置。")
    st.stop()

# 初始化双引擎
client_brain = OpenAI(api_key=api_key_brain, base_url="https://api.deepseek.com")
client_paint = ZhipuAI(api_key=api_key_paint)

if "messages" not in st.session_state:
    st.session_state.messages = [{"role": "assistant", "content": "老师好！我已经升级了**强力文档解析**和**双轨配图系统**。对于创意引入，我会认真作画；对于习题和严谨电路图，我会为您精准留出截图位置！"}]
if "ppt_data" not in st.session_state:
    st.session_state.ppt_data = None
if "current_context" not in st.session_state:
    st.session_state.current_context = ""

# ==========================================
# 2. 强力文档解析引擎 (修复读取失败问题)
# ==========================================
def read_file(file_obj):
    """强力文件读取，兼容多种编码格式"""
    try:
        if file_obj.name.endswith('.txt'):
            content_bytes = file_obj.getvalue()
            try:
                # 尝试用 UTF-8 解码 (Mac/Linux 默认)
                return content_bytes.decode("utf-8")
            except UnicodeDecodeError:
                try:
                    # 尝试用 GBK 解码 (Windows 默认)
                    return content_bytes.decode("gbk")
                except Exception:
                    # 如果都失败，强制忽略错误字符读取
                    return content_bytes.decode("utf-8", errors="ignore")
                    
        elif file_obj.name.endswith('.docx'):
            try:
                return "\n".join([para.text for para in docx.Document(file_obj).paragraphs])
            except Exception as e:
                return f"[Docx解析警告: 文件可能包含无法识别的特殊元素 ({e})]"
    except Exception as e:
        return f"[严重错误: 读取 {file_obj.name} 失败 ({e})]"
    return ""

def get_templates():
    if not os.path.exists("templates"): return []
    return [f for f in os.listdir("templates") if f.endswith(".pptx")]

def search_web(query):
    try:
        results = DDGS().text(query, max_results=3)
        return "\n".join([f"- 【{r['title']}】: {r['body']}" for r in results])
    except Exception:
        return "搜索失败"

# ==========================================
# 3. 双轨配图与生成逻辑
# ==========================================
def generate_physics_image(image_prompt):
    """调用智谱大模型作图 (仅限创意类)"""
    try:
        response = client_paint.images.generations(
            model="cogview-3-plus", 
            prompt=f"一张专业、充满高级感的高中物理教学幻灯片配图。要求具有纪录片风格的真实感。画面内容：{image_prompt}。要求：纯白背景或深邃背景，绝不要出现任何英文字母、公式或乱码文字，纯粹展示物理历史场景或宏观现象。",
        )
        image_url = response.data[0].url
        img_data = requests.get(image_url).content
        return img_data
    except Exception as e:
        return None

def generate_word_document(ppt_data, topic_name):
    doc = docx.Document()
    styles_to_change = ['Normal', 'List Bullet', 'Title', 'Heading 1', 'Heading 2']
    for style_name in styles_to_change:
        try:
            style = doc.styles[style_name]
            style.font.name = '仿宋'
            style._element.rPr.rFonts.set(qn('w:eastAsia'), '仿宋')
        except KeyError:
            pass 
    
    h = doc.add_heading(f"【教学讲义】{topic_name}", 0)
    for run in h.runs:
        run.font.name = '仿宋'
        run._element.rPr.rFonts.set(qn('w:eastAsia'), '仿宋')
    
    for i, slide in enumerate(ppt_data):
        h2 = doc.add_heading(f"环节 {i+1}：{slide.get('title', '无标题')}", level=2)
        for run in h2.runs:
            run.font.name = '仿宋'
            run._element.rPr.rFonts.set(qn('w:eastAsia'), '仿宋')
            
        for point in slide.get("content", []):
            p = doc.add_paragraph(point, style='List Bullet')
            for run in p.runs:
                run.font.name = '仿宋'
                run._element.rPr.rFonts.set(qn('w:eastAsia'), '仿宋')
        
        image_type = slide.get("image_type", "none")
        if image_type == "creative" and slide.get("image_bytes"):
            img_stream = io.BytesIO(slide["image_bytes"])
            doc.add_picture(img_stream, width=docx.shared.Inches(4.0))
            p = doc.add_paragraph("图：AI生成创意配图")
            p.runs[0].font.italic = True
        elif image_type == "schematic":
            p = doc.add_paragraph(f"✂️ [此处预留空位，请手动截取原卷图片粘贴]")
            p.runs[0].font.color.rgb = docx.shared.RGBColor(0, 112, 192)
            
        for p in doc.paragraphs[-2:]: # Ensure recent paragraphs have font set
            for run in p.runs:
                run.font.name = '仿宋'
                run._element.rPr.rFonts.set(qn('w:eastAsia'), '仿宋')
    
    stream = io.BytesIO()
    doc.save(stream)
    stream.seek(0)
    return stream

# ==========================================
# 4. 界面布局
# ==========================================
with st.sidebar:
    st.header("📂 本地资料库")
    uploaded_files = st.file_uploader("多选电脑教案 (按住 Command)：", type=['txt', 'docx'], accept_multiple_files=True)
    if uploaded_files:
        if st.button("📥 强力解析并提取", use_container_width=True):
            with st.spinner("正在解析文档..."):
                all_content = []
                for f in uploaded_files:
                    content = read_file(f)
                    if "[严重错误" in content:
                        st.error(content)
                    else:
                        all_content.append(f"《{f.name}》\n{content}")
                
                if all_content:
                    st.session_state.current_context = "\n\n".join(all_content)
                    st.success("✅ 文档读取完毕！")
                
    st.markdown("---")
    topic = st.text_input("备课课题：", placeholder="例如：法拉第电磁感应定律")
    template_files = get_templates()
    selected_template = st.selectbox("PPT 模板：", template_files) if template_files else None
    selected_template_path = os.path.join("templates", selected_template) if selected_template else None

col_chat, col_studio = st.columns([6, 4], gap="large")

with col_chat:
    st.header("💬 智能教研控制台")
    chat_container = st.container(height=500)
    with chat_container:
        for msg in st.session_state.messages:
            st.chat_message(msg["role"]).write(msg["content"])

    st.markdown("⚡ **Agent 指令面板**")
    btn_col1, btn_col2 = st.columns(2)
    
    with btn_col1:
        if st.button("🕰️ 寻找硬核物理学史", use_container_width=True):
            if not topic: st.warning("先输入课题")
            else:
                st.session_state.messages.append({"role": "user", "content": f"帮我找“{topic}”的真实学史和困境。"})
                with chat_container:
                    st.chat_message("user").write(st.session_state.messages[-1]["content"])
                    with st.chat_message("assistant"):
                        with st.spinner("🌍 翻阅历史文献中..."):
                            search_results = search_web(f"{topic} 物理学史 真实历史事件")
                            prompt = f"你是一位严谨的【科学史学者】。请为“{topic}”写一段硬核的引课，拒绝童话感，必须基于以下真实年份和困境：\n{search_results}"
                            response = client_brain.chat.completions.create(model="deepseek-chat", messages=[{"role": "user", "content": prompt}], temperature=0.3)
                            reply = response.choices[0].message.content
                            st.write(reply)
                            st.session_state.current_context += "\n\n" + reply 
                            st.session_state.messages.append({"role": "assistant", "content": reply})

    with btn_col2:
        if st.button("🪄 双轨智能排版 (生图+预留位)", type="primary", use_container_width=True):
            if not topic: st.warning("先输入课题")
            else:
                st.session_state.messages.append({"role": "user", "content": f"请为“{topic}”生成教案，并智能判断需要 AI 作图还是预留手工截图位。"})
                with chat_container:
                    st.chat_message("user").write(st.session_state.messages[-1]["content"])
                    with st.chat_message("assistant"):
                        with st.spinner("🧠 大脑正在智能分类配图需求..."):
                            # 终极分类 Prompt
                            prompt = f"""
                            参考资料：{st.session_state.current_context}
                            请为课题“{topic}”生成教案大纲。
                            【核心任务】：对于每一页，你需要判断配图的类型（image_type）。
                            - 如果是历史人物、宏观场景、趣味引入等适合 AI 发挥的画面，设为 "creative"，并在 "image_prompt" 中用英文写下高质量画图提示词。
                            - 如果是【具体的习题图】、【受力分析图】、【电路图】、【严格的原理示意图】，设为 "schematic"，并让 "image_prompt" 为空。
                            - 如果不需要图，设为 "none"。
                            
                            【输出格式】必须是纯 JSON 数组：
                            [ 
                                {{
                                    "title": "标题", 
                                    "content": ["要点1"], 
                                    "image_type": "creative", 
                                    "image_prompt": "A cinematic shot of Newton..." 
                                }},
                                {{
                                    "title": "例题解析", 
                                    "content": ["要点1"], 
                                    "image_type": "schematic", 
                                    "image_prompt": "" 
                                }}
                            ]
                            """
                            response = client_brain.chat.completions.create(model="deepseek-chat", messages=[{"role": "user", "content": prompt}], temperature=0.2)
                            result_text = response.choices[0].message.content.replace("```json", "").replace("```", "").strip()
                            try:
                                ppt_data = json.loads(result_text)
                            except Exception as e:
                                st.error("JSON 解析失败，请重试。")
                                st.stop()
                        
                        with st.spinner("🎨 画师正在专心绘制创意插图..."):
                            for slide in ppt_data:
                                if slide.get("image_type") == "creative" and slide.get("image_prompt"):
                                    img_bytes = generate_physics_image(slide["image_prompt"])
                                    if img_bytes:
                                        slide["image_bytes"] = img_bytes
                            
                            st.session_state.ppt_data = ppt_data
                            reply = "✅ 双轨图文教案生成完毕！您可以去右侧 Studio 验收了。"
                            st.write(reply)
                            st.session_state.messages.append({"role": "assistant", "content": reply})

    # ================= 替换从这里开始 =================
    if prompt := st.chat_input("输入指令试试：把文档前三道题做成PPT..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with chat_container:
            st.chat_message("user").write(prompt)
            with st.chat_message("assistant"):
                
                # 【智能路由】：判断用户是想“聊天”还是想“做PPT”
                if "ppt" in prompt.lower() or "幻灯片" in prompt or "排版" in prompt:
                    with st.spinner("🧠 收到排版指令，正在精确提取内容并生成 PPT..."):
                        # 生成 PPT 的专属指令，强制要求带上文档记忆和用户特殊要求
                        ppt_prompt = f"""
                        参考本地文档资料：\n{st.session_state.current_context}\n
                        
                        用户的新指令是：“{prompt}”
                        请严格按照用户的指令（例如只选取前三道题），从资料中提取内容生成 PPT 大纲。
                        
                        【核心任务】：对于每一页，你需要判断配图的类型（image_type）。
                        - 宏观场景/历史人物 -> "creative" (并在 image_prompt 写英文提示词)
                        - 具体的习题图/电路图 -> "schematic" (image_prompt 为空)
                        - 不需要图 -> "none"
                        
                        【输出格式】必须是纯 JSON 数组：
                        [ {{"title": "标题", "content": ["要点1"], "image_type": "schematic", "image_prompt": ""}} ]
                        """
                        try:
                            response = client_brain.chat.completions.create(model="deepseek-chat", messages=[{"role": "user", "content": ppt_prompt}], temperature=0.2)
                            result_text = response.choices[0].message.content.replace("```json", "").replace("```", "").strip()
                            ppt_data = json.loads(result_text)
                            
                            # 创意图 AI 绘制
                            for slide in ppt_data:
                                if slide.get("image_type") == "creative" and slide.get("image_prompt"):
                                    img_bytes = generate_physics_image(slide["image_prompt"])
                                    if img_bytes: slide["image_bytes"] = img_bytes
                            
                            # 将结果保存并推送到右侧 Studio
                            st.session_state.ppt_data = ppt_data
                            reply = "✅ 没问题！我已经按您的要求提取了文档内容，并把生成的 PPT 发送到右侧的 Studio 工作区了，请查收！"
                            st.write(reply)
                            st.session_state.messages.append({"role": "assistant", "content": reply})
                            # 刷新网页渲染右侧界面
                            st.rerun() 
                            
                        except Exception as e:
                            error_msg = f"生成 PPT 失败，请检查文档内容或重试。（报错：{e}）"
                            st.error(error_msg)
                            st.session_state.messages.append({"role": "assistant", "content": error_msg})
                            
                else:
                    # 普通聊天的逻辑（治好“失忆症”）
                    with st.spinner("思考中..."):
                        # 核心修复：把提取的文档内容塞进 AI 的系统大脑里
                        sys_prompt = "你是资深物理老师。"
                        if st.session_state.current_context:
                            sys_prompt += f"\n\n请严格基于以下我提供的本地文档资料来回答问题：\n{st.session_state.current_context}"
                        
                        context_msg = [{"role": "system", "content": sys_prompt}] + st.session_state.messages
                        
                        response = client_brain.chat.completions.create(model="deepseek-chat", messages=context_msg, temperature=0.6)
                        reply = response.choices[0].message.content
                        st.write(reply)
                        st.session_state.messages.append({"role": "assistant", "content": reply})
    # ================= 替换到这里结束 =================

with col_studio:
    st.header("✨ Studio 成果区")
    if st.session_state.ppt_data is None:
        st.info("👈 点击中栏的【双轨智能排版】，查看智能图文分布。")
    else:
        if selected_template_path: prs = Presentation(selected_template_path)
        else: prs = Presentation() 

        for slide_data in st.session_state.ppt_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if slide.shapes.title: slide.shapes.title.text = slide_data.get("title", "无标题")
            
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                contents = slide_data.get("content", [""])
                if contents:
                    tf.text = contents[0]
                    for point in contents[1:]:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
            
            # 双轨配图逻辑
            image_type = slide_data.get("image_type", "none")
            
            if image_type == "creative" and slide_data.get("image_bytes"):
                # 创意图：贴上真实图片
                img_stream = io.BytesIO(slide_data["image_bytes"])
                try:
                    slide.shapes.add_picture(img_stream, Inches(5), Inches(2), width=Inches(4))
                except Exception:
                    pass
            elif image_type == "schematic":
                # 习题图：生成显眼的蓝色预留文字
                try:
                    txBox = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(3))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = "✂️ 【截图预留位】\n请在此处使用快捷键\n粘贴原卷中的习题配图"
                    p.font.color.rgb = RGBColor(0, 112, 192)
                    p.font.size = Pt(20)
                except Exception:
                    pass

        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        word_stream = generate_word_document(st.session_state.ppt_data, topic if topic else "未命名课题")
        
        st.success("🎉 生成成功！")
        dl_col1, dl_col2 = st.columns(2)
        with dl_col1:
            st.download_button("📊 下载 智能双轨 PPT", data=ppt_stream, file_name=f"{topic}_AI备课.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", type="primary", use_container_width=True)
        with dl_col2:
            st.download_button("📝 下载 仿宋 Word 讲义", data=word_stream, file_name=f"{topic}_教案讲义.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", type="secondary", use_container_width=True)
        
        st.markdown("### 👀 图文讲义预览")
        with st.container(height=500):
            for i, slide_data in enumerate(st.session_state.ppt_data):
                with st.container(border=True):
                    st.markdown(f"**环节 {i+1}：{slide_data.get('title', '无标题')}**")
                    for point in slide_data.get("content", []): st.markdown(f"- {point}")
                    
                    if slide_data.get("image_type") == "creative" and slide_data.get("image_bytes"):
                        st.image(slide_data["image_bytes"], caption="AI 绘制的创意配图", width=300)
                    elif slide_data.get("image_type") == "schematic":
                        st.info("✂️ Agent 已为本页预留【原题截图空位】，请下载 PPT 后手动粘贴物理图。")